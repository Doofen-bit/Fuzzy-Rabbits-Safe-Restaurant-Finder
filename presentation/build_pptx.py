"""
build_pptx.py
-------------
Generates the project presentation PowerPoint:
    NYC Safe Restaurant Finder — Machine Learning Showcase

Run from the repo root:
    python presentation/build_pptx.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ── Palette ──────────────────────────────────────────────────────────────────
C_BG       = RGBColor(0x0D, 0x11, 0x17)   # near-black background
C_PANEL    = RGBColor(0x16, 0x21, 0x33)   # dark navy panel
C_ACCENT   = RGBColor(0x4C, 0x8B, 0xF5)   # blue accent
C_GREEN    = RGBColor(0x00, 0xC8, 0x64)   # grade A green
C_YELLOW   = RGBColor(0xFF, 0xC8, 0x00)   # grade B yellow
C_RED      = RGBColor(0xF0, 0x50, 0x3C)   # grade C red
C_TEAL     = RGBColor(0x00, 0xC8, 0xB4)   # teal highlight
C_PURPLE   = RGBColor(0xAA, 0x44, 0xFF)   # purple highlight
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY    = RGBColor(0xCC, 0xCC, 0xCC)   # light gray text
C_DGRAY    = RGBColor(0x66, 0x77, 0x88)   # dim gray
C_ORANGE   = RGBColor(0xFF, 0x8C, 0x00)   # route orange

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helper utilities ──────────────────────────────────────────────────────────

def _prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def _fill_bg(slide, color: RGBColor = C_BG) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(
    slide,
    left: float, top: float, width: float, height: float,
    fill: RGBColor | None = None,
    line_color: RGBColor | None = None,
    line_width_pt: float = 0,
    transparency: float = 0.0,
):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()

    if line_color and line_width_pt > 0:
        shape.line.color.rgb = line_color
        shape.line.width    = Pt(line_width_pt)
    else:
        shape.line.fill.background()

    return shape


def _text_box(
    slide,
    text: str,
    left: float, top: float, width: float, height: float,
    font_size: int    = 18,
    bold:       bool  = False,
    italic:     bool  = False,
    color:      RGBColor = C_WHITE,
    align:      PP_ALIGN = PP_ALIGN.LEFT,
    word_wrap:  bool  = True,
    bg_color:   RGBColor | None = None,
    line_color: RGBColor | None = None,
):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    else:
        txBox.fill.background()

    if line_color:
        txBox.line.color.rgb = line_color
        txBox.line.width     = Pt(1)
    else:
        txBox.line.fill.background()

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text            = text
    run.font.size       = Pt(font_size)
    run.font.bold       = bold
    run.font.italic     = italic
    run.font.color.rgb  = color
    run.font.name       = "Segoe UI"
    return txBox


def _multiline_text(
    slide,
    lines: list[tuple[str, int, bool, RGBColor]],   # (text, size, bold, color)
    left: float, top: float, width: float, height: float,
    bg_color: RGBColor | None = None,
    line_spacing_pt: float = 6,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    else:
        txBox.fill.background()
    txBox.line.fill.background()

    first = True
    for text, size, bold, color in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(line_spacing_pt)
        run = p.add_run()
        run.text           = text
        run.font.size      = Pt(size)
        run.font.bold      = bold
        run.font.color.rgb = color
        run.font.name      = "Segoe UI"
    return txBox


def _pill(
    slide,
    text: str,
    left: float, top: float, width: float, height: float,
    bg: RGBColor = C_ACCENT, fg: RGBColor = C_WHITE,
    font_size: int = 13, bold: bool = True,
):
    """Rounded-corner pill-style label."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        5,   # rounded rectangle
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text           = text
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.color.rgb = fg
    run.font.name      = "Segoe UI"
    return shape


def _accent_bar(slide, top: float = 0.0, height: float = 0.08,
                color: RGBColor = C_ACCENT) -> None:
    _box(slide, 0, top, 13.33, height, fill=color)


def _section_header(slide, number: str, title: str,
                    color: RGBColor = C_ACCENT) -> None:
    """Left colour strip + part number + title."""
    _box(slide, 0, 0, 0.12, 7.5, fill=color)
    _text_box(slide, number, 0.22, 0.18, 2.0, 0.5,
              font_size=13, bold=True, color=color)
    _text_box(slide, title, 0.22, 0.42, 12.5, 0.7,
              font_size=32, bold=True, color=C_WHITE)


# ── Individual slide builders ─────────────────────────────────────────────────

def slide_title(prs: Presentation) -> None:
    slide = _blank(prs)
    _fill_bg(slide)

    # Top accent stripe
    _accent_bar(slide, top=0.0, height=0.55, color=C_ACCENT)

    # Main title
    _text_box(slide, "NYC Safe Restaurant Finder",
              0.5, 0.05, 12.0, 0.55,
              font_size=34, bold=True, color=C_WHITE,
              align=PP_ALIGN.CENTER)

    # Subtitle
    _text_box(slide, "A Machine Learning Showcase — Parts 1 through 6",
              0.5, 1.0, 12.33, 0.5,
              font_size=20, bold=False, color=C_LGRAY,
              align=PP_ALIGN.CENTER)

    # Dataset pill
    _pill(slide, "295,995 DOHMH Inspection Records  ·  30,935 Unique Restaurants  ·  NYC 2020–2026",
          0.5, 1.7, 12.33, 0.42,
          bg=C_PANEL, fg=C_LGRAY, font_size=13, bold=False)

    # Six-part cards
    parts = [
        ("1", "Data &\nExploration",   C_BG,    C_ACCENT),
        ("2", "KNN\nClassifier",       C_BG,    C_GREEN),
        ("3", "Decision\nTree",        C_BG,    C_YELLOW),
        ("4", "Cuisine\nPredictor",    C_BG,    C_TEAL),
        ("5", "RL Route\nFinder",      C_BG,    C_ORANGE),
        ("6", "Ultimate\nFinder",      C_BG,    C_PURPLE),
    ]
    card_w, gap = 1.9, 0.25
    start_x = (13.33 - (card_w * 6 + gap * 5)) / 2

    for i, (num, label, bg, color) in enumerate(parts):
        cx = start_x + i * (card_w + gap)
        _box(slide, cx, 2.45, card_w, 1.6,
             fill=RGBColor(
                 min(0x16 + 0x0A * i, 0xFF),
                 min(0x21 + 0x05 * i, 0xFF),
                 min(0x33 + 0x05 * i, 0xFF),
             ),
             line_color=color, line_width_pt=2)
        _text_box(slide, f"Part {num}", cx + 0.1, 2.5, card_w - 0.2, 0.35,
                  font_size=11, bold=True, color=color,
                  align=PP_ALIGN.CENTER)
        _text_box(slide, label, cx + 0.1, 2.85, card_w - 0.2, 0.9,
                  font_size=14, bold=True, color=C_WHITE,
                  align=PP_ALIGN.CENTER)

    # Bottom "flow" arrow
    _text_box(slide, "Data  →  Model  →  Reinforcement Learning  →  Ultimate Finder",
              0.5, 4.3, 12.33, 0.4,
              font_size=15, color=C_DGRAY, align=PP_ALIGN.CENTER)

    # Bottom bar
    _accent_bar(slide, top=7.3, height=0.2, color=C_PANEL)
    _text_box(slide, "DOHMH NYC Open Data  ·  NumPy · scikit-learn · Folium · Streamlit",
              0.5, 7.28, 12.33, 0.22,
              font_size=10, color=C_DGRAY, align=PP_ALIGN.CENTER)


def slide_dataset(prs: Presentation) -> None:
    slide = _blank(prs)
    _fill_bg(slide)
    _section_header(slide, "Overview", "Dataset & Project Architecture", C_ACCENT)

    # Metric boxes
    metrics = [
        ("295,995",  "violation records"),
        ("30,935",   "unique restaurants"),
        ("27",       "raw CSV columns"),
        ("5 boroughs", "NYC coverage"),
    ]
    mw = 2.7
    mx0 = 0.5
    for i, (val, lbl) in enumerate(metrics):
        x = mx0 + i * (mw + 0.35)
        _box(slide, x, 1.35, mw, 1.1, fill=C_PANEL,
             line_color=C_ACCENT, line_width_pt=1.5)
        _text_box(slide, val, x + 0.1, 1.42, mw - 0.2, 0.52,
                  font_size=26, bold=True, color=C_ACCENT,
                  align=PP_ALIGN.CENTER)
        _text_box(slide, lbl, x + 0.1, 1.92, mw - 0.2, 0.3,
                  font_size=12, color=C_LGRAY, align=PP_ALIGN.CENTER)

    # Pipeline flow (text boxes as diagram)
    steps = [
        ("Raw CSV", "295,995 rows\n27 columns", C_ACCENT),
        ("load_raw()", "Type casting\nNaN handling\nDate parsing", C_GREEN),
        ("build_restaurant_table()", "1 row / restaurant\nAgg features\nGrade encoding", C_YELLOW),
        ("ML Models", "KNN · DT · Cuisine\nCombined · RL", C_PURPLE),
        ("Streamlit", "Interactive\nDashboard", C_ORANGE),
    ]
    sw = 2.1
    sx0 = 0.3
    sy = 2.75
    sh = 1.6

    for i, (title, desc, col) in enumerate(steps):
        sx = sx0 + i * (sw + 0.52)
        _box(slide, sx, sy, sw, sh, fill=C_PANEL,
             line_color=col, line_width_pt=2)
        _text_box(slide, title, sx + 0.08, sy + 0.08, sw - 0.16, 0.38,
                  font_size=12, bold=True, color=col,
                  align=PP_ALIGN.CENTER)
        _text_box(slide, desc, sx + 0.08, sy + 0.5, sw - 0.16, 0.95,
                  font_size=10, color=C_LGRAY, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            _text_box(slide, "→", sx + sw + 0.08, sy + 0.6, 0.4, 0.5,
                      font_size=22, bold=True, color=C_DGRAY,
                      align=PP_ALIGN.CENTER)

    # Key features used
    _text_box(slide, "Key aggregated features per restaurant",
              0.5, 4.55, 12.0, 0.35,
              font_size=13, bold=True, color=C_LGRAY)

    feats = [
        "mean_score  ·  min/max_score  ·  score_range",
        "critical_violations  ·  total_violations  ·  critical_rate",
        "inspection_count  ·  days_since_last_inspection  ·  violations_per_inspection",
        "latest_grade  ·  latitude  ·  longitude  ·  cuisine  ·  boro",
    ]
    for i, f in enumerate(feats):
        _text_box(slide, f"• {f}", 0.7, 4.9 + i * 0.33, 12.0, 0.3,
                  font_size=11, color=C_LGRAY)


def slide_knn(prs: Presentation) -> None:
    slide = _blank(prs)
    _fill_bg(slide)
    _section_header(slide, "Part 2", "KNN Grade Classifier", C_GREEN)

    # Left panel — algorithm
    _box(slide, 0.22, 1.3, 6.3, 5.9, fill=C_PANEL,
         line_color=C_GREEN, line_width_pt=1)
    _text_box(slide, "Algorithm (from scratch — NumPy only)",
              0.35, 1.4, 6.0, 0.35,
              font_size=13, bold=True, color=C_GREEN)

    steps = [
        ("Features (4)", "mean_score · critical_violations\ntotal_violations · days_since_last_inspection"),
        ("Z-score normalise", "μ = train mean  ·  σ = train std\nX̂ = (X − μ) / (σ + ε)"),
        ("Cosine similarity", "sim(q, xᵢ) = (q · xᵢ) / (‖q‖ · ‖xᵢ‖)\nVectorised: (B × F) @ (F × N) matrix"),
        ("Top-K neighbours", "np.argpartition(sims, −K)  →  O(N) partial sort\nK neighbours selected per query"),
        ("Majority vote", "ŷ = argmax  count(grade g among K neighbours)\nOutput: A / B / C"),
    ]
    y = 1.82
    for title, desc in steps:
        _pill(slide, title, 0.35, y, 2.0, 0.3,
              bg=RGBColor(0x16, 0x35, 0x16), fg=C_GREEN, font_size=10, bold=True)
        _text_box(slide, desc, 2.45, y, 4.0, 0.45,
                  font_size=10, color=C_LGRAY)
        y += 0.72

    _text_box(slide, "Temporal split: sorted by latest_inspection_date",
              0.35, 5.5, 6.0, 0.35,
              font_size=10, italic=True, color=C_DGRAY)
    _text_box(slide, "Train = earliest 80%  ·  Test = newest 20%",
              0.35, 5.82, 6.0, 0.3,
              font_size=10, italic=True, color=C_DGRAY)

    # Right panel — the problem
    _box(slide, 6.85, 1.3, 6.25, 5.9, fill=C_PANEL,
         line_color=C_RED, line_width_pt=1)
    _text_box(slide, "Why KNN Fails Here",
              6.98, 1.4, 5.9, 0.35,
              font_size=13, bold=True, color=C_RED)

    problems = [
        (C_RED,    "Class Imbalance",
         "Grade A ≈ 86%  ·  B ≈ 10%  ·  C ≈ 4%\nMajority vote always returns A"),
        (C_RED,    "Lazy Learner",
         "No learned boundary — memorises training\npoints, cannot generalise"),
        (C_RED,    "Cosine Conflation",
         "Angle metric has no causal link to grade\nOutliers dominate after z-scoring"),
        (C_RED,    "No Feature Importance",
         "Cannot tell which features matter\nBlack box — no interpretability"),
        (C_RED,    "O(n) Inference",
         "Every prediction compares against all\n~24 k training restaurants"),
    ]
    y = 1.82
    for col, title, desc in problems:
        _pill(slide, title, 6.98, y, 2.3, 0.3,
              bg=RGBColor(0x33, 0x10, 0x10), fg=col, font_size=10, bold=True)
        _text_box(slide, desc, 9.4, y, 3.55, 0.45,
                  font_size=10, color=C_LGRAY)
        y += 0.72

    _text_box(slide, "→ Part 3 (Decision Tree) addresses all five problems",
              6.98, 5.5, 6.0, 0.55,
              font_size=12, bold=True, color=C_YELLOW)


def slide_dt(prs: Presentation) -> None:
    slide = _blank(prs)
    _fill_bg(slide)
    _section_header(slide, "Part 3", "Decision Tree Grade Classifier", C_YELLOW)

    # Left: features
    _box(slide, 0.22, 1.3, 4.0, 5.9, fill=C_PANEL,
         line_color=C_YELLOW, line_width_pt=1)
    _text_box(slide, "10 Features  (4 base + 6 engineered)",
              0.35, 1.38, 3.75, 0.35,
              font_size=12, bold=True, color=C_YELLOW)

    feat_groups = [
        ("Original (from KNN)", [
            "mean_score",
            "critical_violations",
            "total_violations",
            "days_since_last_inspection",
        ], C_GREEN),
        ("Engineered (new)", [
            "min_score  ·  max_score",
            "score_range = max − min",
            "inspection_count",
            "critical_rate = crit / total",
            "violations_per_inspection",
        ], C_TEAL),
    ]
    y = 1.82
    for group, items, col in feat_groups:
        _text_box(slide, group, 0.35, y, 3.7, 0.28,
                  font_size=11, bold=True, color=col)
        y += 0.3
        for item in items:
            _text_box(slide, f"  • {item}", 0.35, y, 3.7, 0.28,
                      font_size=10, color=C_LGRAY)
            y += 0.28
        y += 0.1

    # Middle: algorithm
    _box(slide, 4.5, 1.3, 4.65, 5.9, fill=C_PANEL,
         line_color=C_YELLOW, line_width_pt=1)
    _text_box(slide, "Weighted Gini Algorithm",
              4.63, 1.38, 4.4, 0.35,
              font_size=12, bold=True, color=C_YELLOW)

    algo = [
        ("Gini Impurity",
         "Gini(S) = 1 − Σ pᵢ²\n(weighted by class frequency)"),
        ("Balanced Class Weights",
         "w(class g) = n / (n_classes · count_g)\nGives B and C equal pull"),
        ("Best Split Selection",
         "Δ = Gini(parent) − Σ wₗ·Gini(left) − wᵣ·Gini(right)\nVectorised over T thresholds × N samples"),
        ("Quantile Thresholds",
         "50 quantile-sampled cut-points per feature\nAvoids evaluating every unique value"),
        ("Stopping Criteria",
         "max_depth | min_samples_split\nmin_samples_leaf | pure node"),
        ("Feature Importance",
         "Σ (weighted Gini decrease at each split)\nNormalised to sum = 1"),
    ]
    y = 1.82
    for title, desc in algo:
        _pill(slide, title, 4.63, y, 2.2, 0.28,
              bg=RGBColor(0x33, 0x2A, 0x00), fg=C_YELLOW, font_size=10, bold=True)
        _text_box(slide, desc, 6.95, y, 2.12, 0.42,
                  font_size=9, color=C_LGRAY)
        y += 0.65

    # Right: comparison table
    _box(slide, 9.42, 1.3, 3.68, 5.9, fill=C_PANEL,
         line_color=C_GREEN, line_width_pt=1)
    _text_box(slide, "DT vs KNN",
              9.55, 1.38, 3.4, 0.35,
              font_size=12, bold=True, color=C_GREEN)

    rows = [
        ("Issue",         "KNN",     "DT Fix"),
        ("Imbalance",     "Fail",    "Balanced weights"),
        ("Boundary",      "None",    "Learned thresholds"),
        ("Metric",        "Cosine",  "Gini impurity"),
        ("Importance",    "None",    "Split Δ Gini"),
        ("Inference",     "O(n)",    "O(depth)"),
    ]
    y = 1.82
    col_colors = [C_LGRAY, C_RED, C_GREEN]
    for ri, row in enumerate(rows):
        for ci, cell in enumerate(row):
            bld = ri == 0
            c   = col_colors[ci] if ri == 0 else (C_LGRAY if ci == 0 else (C_RED if ci == 1 else C_GREEN))
            _text_box(slide, cell,
                      9.55 + ci * 1.2, y,
                      1.18, 0.32,
                      font_size=9, bold=bld, color=c,
                      align=PP_ALIGN.CENTER)
        y += 0.5
        if ri == 0:
            _box(slide, 9.45, y - 0.05, 3.55, 0.03, fill=C_DGRAY)


def slide_cuisine(prs: Presentation) -> None:
    slide = _blank(prs)
    _fill_bg(slide)
    _section_header(slide, "Part 4", "Cuisine Type Predictor", C_TEAL)

    # Pipeline diagram
    pipeline = [
        ("Restaurant\nName (DBA)", "Input", C_LGRAY),
        ("Text\nNormalisation", "lowercase\nstrip punctuation\ncollapse spaces", C_TEAL),
        ("TF-IDF\nVectoriser", "analyzer='char_wb'\nn-gram range (2,4)\nmax 50k features\nsublinear TF", C_ACCENT),
        ("Logistic\nRegression", "C=1.0 · L2 reg\nbalanced weights\nsolver='lbfgs'\nmax_iter=1000", C_GREEN),
        ("Top-3\nPredictions", "Cuisine label\n+ probability\n(softmax)", C_YELLOW),
    ]

    pw = 2.15
    px0 = 0.35
    py = 1.45
    ph = 2.5

    for i, (title, desc, col) in enumerate(pipeline):
        px = px0 + i * (pw + 0.35)
        _box(slide, px, py, pw, ph, fill=C_PANEL,
             line_color=col, line_width_pt=2)
        _text_box(slide, title, px + 0.08, py + 0.1, pw - 0.16, 0.55,
                  font_size=12, bold=True, color=col,
                  align=PP_ALIGN.CENTER)
        _text_box(slide, desc, px + 0.08, py + 0.72, pw - 0.16, 1.5,
                  font_size=10, color=C_LGRAY, align=PP_ALIGN.CENTER)
        if i < len(pipeline) - 1:
            _text_box(slide, "→", px + pw + 0.02, py + ph / 2 - 0.2, 0.35, 0.4,
                      font_size=20, bold=True, color=C_DGRAY,
                      align=PP_ALIGN.CENTER)

    # Why char n-grams
    _box(slide, 0.22, 4.2, 6.2, 2.9, fill=C_PANEL,
         line_color=C_TEAL, line_width_pt=1)
    _text_box(slide, "Why Character N-grams?",
              0.4, 4.3, 5.9, 0.35,
              font_size=12, bold=True, color=C_TEAL)
    _multiline_text(slide, [
        ("• Capture sub-word patterns: 'pizz' → pizza · 'sush' → sushi", 10, False, C_LGRAY),
        ("• Language-agnostic: works across cuisines without tokenisation", 10, False, C_LGRAY),
        ("• Robust to spelling variants and abbreviations", 10, False, C_LGRAY),
        ("• char_wb pads each word boundary with whitespace markers", 10, False, C_LGRAY),
        ("• n-gram (2,4): bigrams + trigrams + quadgrams = rich vocabulary", 10, False, C_LGRAY),
    ], 0.4, 4.7, 5.8, 2.1, line_spacing_pt=4)

    # Split methods
    _box(slide, 6.72, 4.2, 6.38, 2.9, fill=C_PANEL,
         line_color=C_GREEN, line_width_pt=1)
    _text_box(slide, "Two Train/Test Split Strategies",
              6.9, 4.3, 6.0, 0.35,
              font_size=12, bold=True, color=C_GREEN)
    _multiline_text(slide, [
        ("Random split (stratified):", 11, True, C_GREEN),
        ("  Randomly stratify 80/20 across all boroughs", 10, False, C_LGRAY),
        ("  Tests generalisation across all cuisine types", 10, False, C_LGRAY),
        ("Geographic (borough hold-out):", 11, True, C_TEAL),
        ("  One entire NYC borough excluded from training", 10, False, C_LGRAY),
        ("  Tests spatial generalisation — does the model", 10, False, C_LGRAY),
        ("  predict cuisine from restaurant names it's never seen?", 10, False, C_LGRAY),
    ], 6.9, 4.72, 6.1, 2.1, line_spacing_pt=3)


def slide_rl(prs: Presentation) -> None:
    slide = _blank(prs)
    _fill_bg(slide)
    _section_header(slide, "Part 5", "RL-Powered Safe Restaurant Route Finder", C_ORANGE)

    # MDP components
    _box(slide, 0.22, 1.3, 5.5, 5.9, fill=C_PANEL,
         line_color=C_ORANGE, line_width_pt=1)
    _text_box(slide, "MDP Formulation",
              0.35, 1.38, 5.2, 0.35,
              font_size=12, bold=True, color=C_ORANGE)

    mdp = [
        ("State S",
         "Grid cell (row, col) ≈ 400 m × 420 m\n111 × 113 = 12,543 cells covering all NYC"),
        ("Action A",
         "Move in 8 compass directions:\nN · NE · E · SE · S · SW · W · NW"),
        ("Reward R(s)",
         "Safety scores (A=3 B=2 C=1) of cuisine-filtered\nrestaurants in cell s, Gaussian-blurred\n× cuisine TF-IDF match score"),
        ("Discount γ",
         "Derived from walking-time budget:\nγ ∈ {0.50 … 0.85} (5–20 min presets)"),
        ("Algorithm",
         "Value Iteration — Bellman update (NumPy):\nVⁿ⁺¹(s) = R(s) + γ · max_a Vⁿ(T(s,a))\nFully vectorised via np.pad + np.maximum.reduce"),
    ]
    y = 1.82
    for title, desc in mdp:
        _pill(slide, title, 0.35, y, 1.6, 0.28,
              bg=RGBColor(0x33, 0x1C, 0x00), fg=C_ORANGE, font_size=10)
        _text_box(slide, desc, 2.1, y, 3.5, 0.5,
                  font_size=10, color=C_LGRAY)
        y += 0.72

    # Right: pipeline steps
    _box(slide, 6.0, 1.3, 7.1, 5.9, fill=C_PANEL,
         line_color=C_ORANGE, line_width_pt=1)
    _text_box(slide, "Route Planning Pipeline",
              6.15, 1.38, 6.8, 0.35,
              font_size=12, bold=True, color=C_ORANGE)

    pipe = [
        ("1  Cuisine NLP match",
         "TF-IDF char n-grams on cuisine labels\nExpand food keywords: 'ramen' → 'Japanese noodles'\nCosine similarity → per-cuisine match ∈ [0.05, 1.0]"),
        ("2  Build reward grid",
         "For each restaurant: weighted_score = grade_safety × cuisine_match\nnp.add.at accumulates into (111×113) grid\ngaussian_filter(σ) spreads density to adjacent cells"),
        ("3  Proximity Gaussian bias",
         "R_biased = R × exp(−dist² / 2σ²)\nSuppresses far clusters beyond walking radius\nσ derived from walk-time preset"),
        ("4  Value Iteration",
         "Bellman update converges in <400 iterations\nAll 8-direction neighbours via np.pad shifts\nConvergence check: max|Vⁿ⁺¹ − Vⁿ| < 1e-5"),
        ("5  Greedy trace + OSRM",
         "Follow ∇V from start → destination cell\nReal pedestrian route via OSRM API (walking mode)\nFallback: rectilinear street-grid approximation"),
    ]
    y = 1.82
    for i, (title, desc) in enumerate(pipe):
        _pill(slide, title, 6.15, y, 2.5, 0.28,
              bg=RGBColor(0x33, 0x1C, 0x00), fg=C_ORANGE, font_size=10)
        _text_box(slide, desc, 6.15, y + 0.32, 6.7, 0.5,
                  font_size=9.5, color=C_LGRAY)
        y += 1.05


def slide_ultimate(prs: Presentation) -> None:
    slide = _blank(prs)
    _fill_bg(slide)
    _section_header(slide, "Part 6", "Ultimate Restaurant Finder — Everything Combined", C_PURPLE)

    # Three panels
    panels = [
        (0.22,  3.7, C_PURPLE, "Combined Grade Predictor",
         [
             ("Why blend?", "KNN skews → A  (class imbalance, no weighting)\nDT skews → B/C  (balanced Gini)\nBlend gives calibrated continuous score", C_LGRAY),
             ("Probability blend", "P_combined = α·P_KNN + (1−α)·P_DT\n(default α = 0.35)", C_PURPLE),
             ("Safety score", "score = 3·P(A) + 2·P(B) + 1·P(C) ∈ [1,3]\nReplaces binary A→3 / B→2 / C→1 look-up", C_PURPLE),
             ("KNN proba", "Fraction of each grade among K neighbours\nChunked (256/batch) for RAM efficiency", C_LGRAY),
             ("DT proba", "predict_proba() from weighted Gini tree\nProbabilities aligned to [A, B, C] order", C_LGRAY),
         ]),
        (4.18,  4.5, C_TEAL,   "Restaurant-Level NLP Embedding",
         [
             ("Corpus", "TF-IDF matrix (char n-grams 2–4)\nOne document per restaurant:\n\"{name} {cuisine} {borough}\"", C_LGRAY),
             ("Query embedding", "User text → vectorizer.transform() → (1×d)\nCosine vs restaurant matrix → (n,) scores", C_TEAL),
             ("Food keywords", "Expansion: 'ramen' → 'japanese noodles asian'\nBridges dish names ↔ cuisine labels", C_LGRAY),
             ("Restaurant reference", "Detect restaurant name in query\n'like Nobu' → use Nobu's own vector as query\nFinds structurally similar restaurants", C_TEAL),
         ]),
        (8.8,   4.2, C_GREEN,  "Dual Navigation Modes",
         [
             ("Area Finder (RL)", "Value Iteration with combined scores:\nreward = safety_score × nlp_score\nRoutes to densest matching cluster", C_LGRAY),
             ("Direct Route", "Rank by safety × nlp × Gaussian_proximity\nRoute straight to top-ranked restaurant\nNext / Previous / Jump-to-# navigation", C_GREEN),
             ("Both modes", "OSRM walking-mode pedestrian routing\nGrade-C warning overlay on map\nRL value-function heatmap visualisation", C_LGRAY),
         ]),
    ]

    for px, pw, col, title, items in panels:
        _box(slide, px, 1.3, pw, 5.85, fill=C_PANEL,
             line_color=col, line_width_pt=1.5)
        _text_box(slide, title, px + 0.12, 1.38, pw - 0.24, 0.38,
                  font_size=12, bold=True, color=col)
        y = 1.85
        for ititle, idesc, ic in items:
            _pill(slide, ititle, px + 0.12, y, pw - 0.24, 0.26,
                  bg=RGBColor(0x18, 0x18, 0x30), fg=col, font_size=9, bold=True)
            _text_box(slide, idesc, px + 0.12, y + 0.3, pw - 0.24, 0.65,
                      font_size=9, color=ic)
            y += 1.02


def slide_combined_detail(prs: Presentation) -> None:
    """Deep-dive on the blended probability math."""
    slide = _blank(prs)
    _fill_bg(slide)
    _section_header(slide, "Part 6 Deep-Dive", "Blended Probability Math & NLP Matrix Embedding", C_PURPLE)

    # Left: math
    _box(slide, 0.22, 1.3, 6.3, 5.9, fill=C_PANEL,
         line_color=C_PURPLE, line_width_pt=1)
    _text_box(slide, "Grade Probability Blending",
              0.35, 1.38, 6.0, 0.35,
              font_size=13, bold=True, color=C_PURPLE)

    _multiline_text(slide, [
        ("Step 1 — KNN Class Probabilities", 11, True, C_GREEN),
        ("  For each query x, find K nearest neighbours:", 10, False, C_LGRAY),
        ("  P_KNN(g | x) = count(grade = g in K neighbours) / K", 10, False, C_TEAL),
        ("  Cosine similarity computed as (B×F)@(F×N) matrix multiply", 10, False, C_LGRAY),
        ("  → chunked in batches of 256 to bound peak RAM", 10, False, C_DGRAY),
        ("", 8, False, C_LGRAY),
        ("Step 2 — Decision Tree Class Probabilities", 11, True, C_YELLOW),
        ("  Leaf node stores weighted fraction per class:", 10, False, C_LGRAY),
        ("  P_DT(g | x) = Σ w[i∈leaf, grade=g] / Σ w[i∈leaf]", 10, False, C_TEAL),
        ("  Weights from balanced class scheme: n / (K · count_g)", 10, False, C_LGRAY),
        ("", 8, False, C_LGRAY),
        ("Step 3 — Weighted Blend", 11, True, C_PURPLE),
        ("  P_combined(g | x) = α·P_KNN(g) + (1−α)·P_DT(g)", 10, False, C_TEAL),
        ("  α = 0.35 (KNN)  ·  (1−α) = 0.65 (DT)  — tunable", 10, False, C_LGRAY),
        ("", 8, False, C_LGRAY),
        ("Step 4 — Continuous Safety Score", 11, True, C_ORANGE),
        ("  safety(x) = 3·P(A|x) + 2·P(B|x) + 1·P(C|x)  ∈ [1, 3]", 10, False, C_TEAL),
        ("  Replaces binary look-up: smoother RL reward signal", 10, False, C_LGRAY),
    ], 0.35, 1.82, 6.0, 5.1, line_spacing_pt=2)

    # Right: NLP embedding pipeline
    _box(slide, 6.82, 1.3, 6.28, 5.9, fill=C_PANEL,
         line_color=C_TEAL, line_width_pt=1)
    _text_box(slide, "NLP Matrix Embedding Pipeline",
              6.95, 1.38, 5.95, 0.35,
              font_size=13, bold=True, color=C_TEAL)

    _multiline_text(slide, [
        ("Building the Restaurant Matrix", 11, True, C_TEAL),
        ("  doc_i = \"{name_i} {cuisine_i} {borough_i}\"", 10, False, C_LGRAY),
        ("  TF-IDF: analyzer='char_wb', n-gram (2,4), max 100k features", 10, False, C_LGRAY),
        ("  R = vectorizer.fit_transform(docs)   shape: (N_rest × d)", 10, False, C_TEAL),
        ("", 8, False, C_LGRAY),
        ("Embedding a User Query", 11, True, C_TEAL),
        ("  Detect restaurant name in query?", 10, False, C_LGRAY),
        ("    YES → q_vec = R[matched_restaurant_idx]  (1×d sparse row)", 10, False, C_GREEN),
        ("    NO  → expand food keywords, then:", 10, False, C_YELLOW),
        ("          q_vec = vectorizer.transform([expanded_query])", 10, False, C_YELLOW),
        ("", 8, False, C_LGRAY),
        ("Per-Restaurant Match Scores", 11, True, C_TEAL),
        ("  sims = cosine_similarity(q_vec, R)   shape: (N_rest,)", 10, False, C_LGRAY),
        ("  nlp_score_i = 0.05 + 0.95·(sims_i / max(sims))", 10, False, C_TEAL),
        ("", 8, False, C_LGRAY),
        ("Integration into RL Reward Grid", 11, True, C_ORANGE),
        ("  weighted_score_i = safety_i × nlp_score_i", 10, False, C_LGRAY),
        ("  np.add.at(R_grid, (row_i, col_i), weighted_score_i)", 10, False, C_LGRAY),
        ("  R_biased = R_grid × Gaussian_proximity(start_cell)", 10, False, C_LGRAY),
        ("  V ← Value Iteration(R_biased, γ)  →  greedy route", 10, False, C_TEAL),
    ], 6.95, 1.82, 6.0, 5.1, line_spacing_pt=2)


def slide_architecture(prs: Presentation) -> None:
    slide = _blank(prs)
    _fill_bg(slide)
    _section_header(slide, "Architecture", "System Design — How the Parts Connect", C_ACCENT)

    # Data flow from left to right
    layers = [
        ("Raw Data", ["DOHMH CSV\n295,995 rows"], C_ACCENT),
        ("Preprocessing", ["data_loader.load_raw()\nType casting · NaN cleaning",
                           "preprocessor.build_\nrestaurant_table()\n1 row / CAMIS"], C_GREEN),
        ("ML Models", ["KNNClassifier\n(cosine, NumPy only)",
                       "DecisionTreeClassifier\n(weighted Gini, NumPy)",
                       "CuisinePredictor\n(TF-IDF + Logistic Reg)"], C_YELLOW),
        ("Part 6 Fusion", ["CombinedGradePredictor\n(blended KNN+DT proba)",
                           "RestaurantEmbedder\n(restaurant matrix NLP)"], C_PURPLE),
        ("RL Routing", ["MealMatcher / Embedder\n(cuisine/restaurant NLP)",
                        "Value Iteration\n(Bellman, NumPy)",
                        "OSRM walking API\n(real street routes)"], C_ORANGE),
        ("Streamlit", ["Interactive dashboard\n6 tabs · Folium maps\nPlotly charts"], C_TEAL),
    ]

    col_w  = 2.0
    col_gap= 0.12
    x0     = 0.15

    for ci, (lname, items, col) in enumerate(layers):
        cx = x0 + ci * (col_w + col_gap)
        # Column header
        _pill(slide, lname, cx, 1.25, col_w, 0.35,
              bg=col, fg=C_BG, font_size=10, bold=True)

        # Cards in column
        item_h = (5.7 - 0.1 * (len(items) - 1)) / len(items)
        for ri, item in enumerate(items):
            ry = 1.72 + ri * (item_h + 0.1)
            _box(slide, cx, ry, col_w, item_h, fill=C_PANEL,
                 line_color=col, line_width_pt=1)
            _text_box(slide, item, cx + 0.08, ry + 0.08, col_w - 0.16, item_h - 0.16,
                      font_size=9, color=C_LGRAY, align=PP_ALIGN.CENTER)

        # Arrow between columns
        if ci < len(layers) - 1:
            ax = cx + col_w + 0.01
            _text_box(slide, "→", ax, 3.5, col_gap + 0.08, 0.4,
                      font_size=16, bold=True, color=C_DGRAY,
                      align=PP_ALIGN.CENTER)

    # Source note
    _text_box(slide,
              "All core ML algorithms (KNN, Decision Tree) implemented from scratch with NumPy only.  "
              "scikit-learn used for TF-IDF, Logistic Regression, and cosine similarity utilities.",
              0.3, 7.08, 12.7, 0.32,
              font_size=10, italic=True, color=C_DGRAY,
              align=PP_ALIGN.CENTER)


def slide_conclusion(prs: Presentation) -> None:
    slide = _blank(prs)
    _fill_bg(slide)
    _accent_bar(slide, top=0.0, height=0.08, color=C_ACCENT)
    _text_box(slide, "Key Takeaways", 0.5, 0.2, 12.0, 0.65,
              font_size=30, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    takeaways = [
        (C_GREEN,  "Parts 1–3 — From Data to Classifiers",
         "Temporal train/test split prevents data leakage.  "
         "Class imbalance is the central challenge: KNN defaults to Grade A while "
         "Decision Tree with balanced Gini weights recovers B/C predictions.  "
         "10 engineered features give the DT richer signal than KNN's 4."),
        (C_TEAL,   "Part 4 — Character N-gram NLP",
         "TF-IDF on character n-grams (2–4) captures sub-word culinary patterns "
         "('sush', 'pizz', 'taco') without any external tokeniser.  "
         "Geographic hold-out split tests whether cuisine knowledge transfers across NYC boroughs."),
        (C_ORANGE, "Part 5 — Reinforcement Learning for Navigation",
         "Value Iteration (pure NumPy) on a 111×113 grid covers all of NYC.  "
         "Proximity Gaussian bias keeps routes realistic.  "
         "OSRM walking-mode API provides real pedestrian routes."),
        (C_PURPLE, "Part 6 — Combining Everything",
         "Blending KNN and DT probability distributions gives a calibrated continuous safety score "
         "that smooths the RL reward landscape.  "
         "Embedding user queries against the full restaurant TF-IDF matrix enables "
         "dish, cuisine, AND similar-restaurant queries in one unified NLP step.  "
         "Dual navigation modes (cluster RL vs direct route) suit different user needs."),
    ]

    y = 1.1
    for col, title, body in takeaways:
        _box(slide, 0.3, y, 12.73, 1.3, fill=C_PANEL,
             line_color=col, line_width_pt=2)
        _text_box(slide, title, 0.5, y + 0.05, 12.3, 0.36,
                  font_size=12, bold=True, color=col)
        _text_box(slide, body, 0.5, y + 0.42, 12.3, 0.78,
                  font_size=10, color=C_LGRAY)
        y += 1.45

    _accent_bar(slide, top=7.3, height=0.2, color=C_PANEL)
    _text_box(slide,
              "NumPy · scikit-learn · Streamlit · Folium · OSRM  ·  "
              "DOHMH NYC Open Data",
              0.5, 7.28, 12.33, 0.22,
              font_size=10, color=C_DGRAY, align=PP_ALIGN.CENTER)


# ── Build ─────────────────────────────────────────────────────────────────────

def build(output: str = "presentation/NYC_Safe_Restaurant_Finder.pptx") -> None:
    prs = _prs()

    slide_title(prs)
    slide_dataset(prs)
    slide_knn(prs)
    slide_dt(prs)
    slide_cuisine(prs)
    slide_rl(prs)
    slide_ultimate(prs)
    slide_combined_detail(prs)
    slide_architecture(prs)
    slide_conclusion(prs)

    out = Path(output)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out}  ({out.stat().st_size // 1024} KB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
